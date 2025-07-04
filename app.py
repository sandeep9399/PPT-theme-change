import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import io

THEME_FILE = "apollo_theme_with_logo_footer.pptx"

def copy_shapes(source_slide, dest_slide):
    for shape in source_slide.shapes:
        try:
            if shape.shape_type == 1 and shape.has_text_frame:
                new_shape = dest_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_shape.text = shape.text
            elif shape.shape_type == 13:
                image_stream = io.BytesIO(shape.image.blob)
                dest_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)
            elif shape.shape_type == 19:  # table
                table = shape.table
                new_table_shape = dest_slide.shapes.add_table(
                    table.rows.__len__(), table.columns.__len__(),
                    shape.left, shape.top, shape.width, shape.height
                )
                for r in range(len(table.rows)):
                    for c in range(len(table.columns)):
                        new_table_shape.table.cell(r, c).text = table.cell(r, c).text
        except Exception:
            continue

def clear_all_slides(ppt):
    slide_ids = list(ppt.slides._sldIdLst)
    for slide_id in slide_ids:
        ppt.slides._sldIdLst.remove(slide_id)

def apply_apollo_theme(uploaded_pptx):
    source_ppt = Presentation(uploaded_pptx)
    theme_ppt = Presentation(THEME_FILE)
    layouts = theme_ppt.slide_layouts
    output_ppt = Presentation(THEME_FILE)
    clear_all_slides(output_ppt)

    for slide in source_ppt.slides:
        layout = layouts[1]
        new_slide = output_ppt.slides.add_slide(layout)
        copy_shapes(slide, new_slide)

    output = io.BytesIO()
    output_ppt.save(output)
    output.seek(0)
    return output

st.set_page_config(page_title="Apollo PPT Themer", layout="centered")
st.title("🎓 Apollo PPT Themer")

uploaded_file = st.file_uploader("📤 Upload a .pptx file", type=["pptx"])

if uploaded_file:
    st.success("Uploaded. Rebuilding with Apollo theme...")

    with st.spinner("Processing..."):
        result = apply_apollo_theme(uploaded_file)

    base_name = uploaded_file.name.replace(".pptx", "")
    st.download_button(
        label="📥 Download Themed PPTX",
        data=result,
        file_name=f"Apollo_Themed_{base_name}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
